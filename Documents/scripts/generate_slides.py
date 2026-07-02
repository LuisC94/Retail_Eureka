from pptx import Presentation
from pptx.util import Inches, Pt

def create_presentation():
    prs = Presentation()

    # --- Slide Title ---
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Integração Blockchain:\nPlataforma Retail-Eureka"
    subtitle.text = "Estratégia de Adaptação e Arquitetura Híbrida"

    # --- Slide 1: Adaptações ao Código (Chaincode) ---
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = "Adaptação da Lógica de Negócio (Chaincode)"

    tf = body_shape.text_frame
    tf.text = "Simplificação de Estados de Transporte:"
    p = tf.add_paragraph()
    p.text = "Substituição da verificação rígida 'Kargolandi' por estados flexíveis (IN_TRANSIT)."
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Remoção de Dados Obrigatórios:"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Eliminação da obrigatoriedade do campo Matrícula (trackingNo) no Pickup."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Foco na rastreabilidade do lote, agilizando o processo."
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Nova Funcionalidade de Tracking:"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Implementação da função GetHistoryForAsset(id)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Permite consultar toda a linha temporal do produto para gerar gráficos."
    p.level = 1

    # --- Slide 2: Arquitetura de Integração ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = "Arquitetura de Integração: Do Web 2.0 para Web 3.0"

    tf = body_shape.text_frame
    tf.text = "Evolução da Arquitetura (Original vs. Nova):"
    p = tf.add_paragraph()
    p.text = "Original: Angular -> .NET -> Middleware Go -> Blockchain."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Nova: Django (Python) -> Middleware Go -> Blockchain."
    p.level = 1

    p = tf.add_paragraph()
    p.text = "O Papel do Middleware ('O Tradutor'):"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "A Blockchain não 'fala' JSON, exige protocolos seguros (gRPC)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "O Django envia JSON simples. O Middleware traduz, assina e envia para a rede."
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Vantagem da Abordagem:"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Isolamento: Mantemos o Django simples (Windows) e a segurança complexa no Linux."
    p.level = 1

    prs.save('Integration_Slide_Deck.pptx')
    print("Presentation saved successfully: Integration_Slide_Deck.pptx")

if __name__ == "__main__":
    create_presentation()
